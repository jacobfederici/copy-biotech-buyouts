#!/usr/bin/env python3
"""Build a per-company buyout deck: timeline, rumors, stock price, events.

    .venv/bin/python buyouts/deck.py LOXO
    .venv/bin/python buyouts/deck.py icagen_inc_ICGN.md
    .venv/bin/python buyouts/deck.py --all          # every report

Reads:  buyouts/reports/<file>.md + buyouts/data/{deals.json, prices.csv}
Writes: buyouts/data/decks/<ticker>.pptx  (native PPTX shapes/charts)
"""
import json
import re
import sys
import datetime as dt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE

from common import (REPORTS, DECKS, DEALS_JSON, load_prices, fnum,
                    get_field, section, parse_table_rows, first_date, num)

# ---- palette (cohesive with the dashboard) ----
NAVY = "#0E1A2B"; PANEL = "#16243A"; INK = "#1B2733"; MUT = "#6B7C8F"
BLUE = "#2E78D2"; HOT = "#E8455F"; TEAL = "#13A88A"; GOLD = "#E89B2B"
LIGHT = "#FFFFFF"; GRID = "#E4EAF1"
def rgb(h): return RGBColor.from_string(h.lstrip("#"))

CATALYST_RE = re.compile(
    r"\b(FDA|approv|breakthrough|fast track|orphan|priority review|"
    r"phase\s*[123]|topline|top-line|pivotal|data|results|"
    r"acqui|merger|definitive agreement|buyout|licen|partner|collaborat)\b", re.I)
EVENT_TAGS = [
    ("Deal/M&A", re.compile(r"acqui|merger|definitive agreement|buyout|tender offer", re.I)),
    ("Regulatory", re.compile(r"FDA|approv|breakthrough|fast track|orphan|priority|NDA|BLA|sNDA|CRL", re.I)),
    ("Clinical", re.compile(r"phase\s*[123]|topline|top-line|pivotal|data|results|trial|enroll|endpoint", re.I)),
    ("Partnership", re.compile(r"licen|partner|collaborat", re.I)),
    ("Financing", re.compile(r"offering|priced|financ|underwrit", re.I)),
]
def tag_event(h):
    for name, rx in EVENT_TAGS:
        if rx.search(h):
            return name
    return "Other"


# ---- faithful date parsing: keep the date AND its precision (deck-specific) ----
_QUAL = re.compile(r"\s*\((?:month|early|late|mid|morning|afternoon|evening|"
                   r"approx\.?|approximately|estimated|est\.?|circa|SC ?14D9)\)", re.I)


def parse_when(raw):
    """Parse a timeline date token of any precision.

    Returns (date, precision) where precision is 'exact'|'month'|'year', or
    (None, None) if unparseable. For ranges, the START is used.
    """
    s = _QUAL.sub("", raw).strip()
    m = re.match(r"(\d{4})-(\d{2})-(\d{2})", s)
    if m:
        y, mo, da = int(m[1]), int(m[2]), int(m[3])
        if 1 <= da <= 31:
            try:
                return dt.date(y, mo, da), "exact"
            except ValueError:
                pass
        try:
            return dt.date(y, mo, 1), "month"
        except ValueError:
            return None, None
    m = re.match(r"(\d{4})-(\d{2})(?!\d)", s)
    if m:
        try:
            return dt.date(int(m[1]), int(m[2]), 1), "month"
        except ValueError:
            return None, None
    m = re.match(r"(\d{4})\b", s)
    if m:
        return dt.date(int(m[1]), 1, 1), "year"
    return None, None


def fmt_when(d, prec):
    if prec == "year":
        return str(d.year)
    if prec == "month":
        return d.strftime("%b %Y")
    return d.strftime("%b %d, %Y")


def resolve_file(arg):
    if arg.endswith(".md"):
        p = REPORTS / arg
        return p if p.exists() else None
    arg_u = arg.upper()
    for p in REPORTS.glob(f"*_{arg_u}.md"):
        return p
    for p in REPORTS.glob("*.md"):
        if p.stem.upper().endswith("_" + arg_u):
            return p
    return None


def parse_report(fp):
    text = fp.read_text(encoding="utf-8", errors="replace")
    d = {"file": fp.name, "title": text.splitlines()[0].lstrip("# ").strip()}
    for k in ("Ticker", "Status"):
        d[k.lower()] = get_field(text, k)
    d["ann"] = first_date(get_field(text, "Announcement date"))
    d["close"] = first_date(get_field(text, "Close date"))
    tl = []
    for line in section(text, "Timeline").splitlines():
        line = line.strip()
        if not line.startswith("-"):
            continue
        m = re.match(r"-\s*(.+?):\s+(.+)", line)
        if not m:
            continue
        when, prec = parse_when(m.group(1))
        if when is None:
            continue
        tl.append({"date": when, "prec": prec, "text": m.group(2).strip()})
    tl.sort(key=lambda x: x["date"])
    d["timeline"] = tl
    news = []
    for cells in parse_table_rows(section(text, "Company News")):
        if len(cells) < 2:
            continue
        dd = first_date(cells[0])
        if dd:
            news.append((dd, cells[1]))
    d["news"] = sorted(news)
    infl = []
    for cells in parse_table_rows(section(text, "Value Inflection Points")):
        if len(cells) < 6:
            continue
        wk = first_date(cells[0])
        if not wk:
            continue
        infl.append({"week": wk, "trigger": cells[1],
                     "pre": num(cells[3]), "post": num(cells[4]),
                     "chg": num(cells[5].replace("%", ""))})
    d["inflections"] = infl
    return d


def load_dataset():
    deals = {x["file"]: x for x in json.loads(DEALS_JSON.read_text())}
    return deals, load_prices()   # prices: file -> [(date, close), ...]


def clip(s, n):
    s = (s or "").strip()
    if len(s) <= n:
        return s
    return s[:n].rsplit(" ", 1)[0].rstrip(",;:") + "…"


def classify_inflections(rep, cap=12):
    """{week_date: 'deal'|'rumor'|'event'} for the most material moves."""
    ann, news = rep["ann"], rep["news"]
    cand = []
    for inf in rep["inflections"]:
        if inf["chg"] is None:
            continue
        wk = inf["week"]
        is_deal = bool(ann and -3 <= (wk - ann).days <= 10)
        near = [h for (nd, h) in news if abs((nd - wk).days) <= 12 and CATALYST_RE.search(h)]
        if is_deal:
            kind = "deal"
        elif ("not identified" in (inf["trigger"] or "").lower()) and not near:
            kind = "rumor"
        else:
            kind = "event"
        cand.append((wk, kind, abs(inf["chg"])))
    deals = [(wk, k) for wk, k, _ in cand if k == "deal"]
    rest = sorted([(wk, k, a) for wk, k, a in cand if k != "deal"],
                  key=lambda t: t[2], reverse=True)[:cap]
    return {wk: k for wk, k in deals} | {wk: k for wk, k, _ in rest}


def add_price_chart(slide, rep, deal, series):
    """Native line chart: weekly close + classified deal/rumor/event markers.

    series is a list of (date, close). Rumor/event markers are shifted one week
    earlier (reported moves can lag the true move); the deal marker stays put.
    """
    LAG_WEEKS = 1
    n = len(series)
    step = max(1, n // 11)
    cats = [(d.strftime("%b '%y") if i % step == 0 else "") for i, (d, _) in enumerate(series)]
    closes = [round(c, 2) for _, c in series]
    cls = classify_inflections(rep)
    idx = {d: i for i, (d, _) in enumerate(series)}
    marker_series = {"deal": [None] * n, "rumor": [None] * n, "event": [None] * n}
    for d, kind in cls.items():
        i = idx.get(d)
        if i is None:
            continue
        tgt = i if kind == "deal" else max(0, i - LAG_WEEKS)
        marker_series[kind][tgt] = round(series[tgt][1], 2)

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Close", closes)
    order = [k for k in ("event", "rumor", "deal") if any(v is not None for v in marker_series[k])]
    for k in order:
        cd.add_series(k.capitalize(), marker_series[k])

    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                                Inches(0.5), Inches(1.45), Inches(12.3), Inches(4.75), cd)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    plot = chart.plots[0]
    s0 = plot.series[0]
    s0.format.line.color.rgb = rgb(BLUE); s0.format.line.width = Pt(2.0)
    s0.smooth = False
    s0.marker.style = XL_MARKER_STYLE.NONE
    colmap = {"Deal": GOLD, "Rumor": HOT, "Event": TEAL}
    for i, ser in enumerate(plot.series):
        if i == 0:
            continue
        ser.format.line.fill.background()
        ser.smooth = False
        ser.marker.style = XL_MARKER_STYLE.CIRCLE
        ser.marker.size = 10
        mf = ser.marker.format
        mf.fill.solid(); mf.fill.fore_color.rgb = rgb(colmap.get(ser.name, MUT))
        mf.line.color.rgb = rgb("FFFFFF"); mf.line.width = Pt(1.0)
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(8); ca.tick_labels.font.color.rgb = rgb(MUT)
    ca.has_major_gridlines = False
    va = chart.value_axis
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = rgb(MUT)
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = rgb(GRID)
    va.major_gridlines.format.line.width = Pt(0.5)


def add_timeline_diagram(slide, rep, left, top, width):
    """One dot per milestone on a true-date line; announcement & close marked."""
    ms = rep["timeline"]
    if not ms:
        return
    ann, close = rep["ann"], rep["close"]
    ds = [m["date"] for m in ms] + ([ann] if ann else []) + ([close] if close else [])
    lo, hi = min(ds), max(ds)
    span = (hi - lo).days or 1
    def X(d):
        return left + (d - lo).days / span * width
    axis_y = top + 0.55
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(left), Inches(axis_y),
                                    Inches(left + width), Inches(axis_y))
    ln.line.color.rgb = rgb(GRID); ln.line.width = Pt(2)
    for d, col, lab, dy in [(ann, GOLD, "Announcement", -0.42), (close, BLUE, "Close", 0.30)]:
        if not d:
            continue
        g = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(X(d)), Inches(top + 0.2),
                                       Inches(X(d)), Inches(axis_y + 0.35))
        g.line.color.rgb = rgb(col); g.line.width = Pt(1.5)
        if col == GOLD:
            g.line.dash_style = None
        tb = slide.shapes.add_textbox(Inches(X(d) - 1.0), Inches(axis_y + dy), Inches(2.0), Inches(0.25))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = lab; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = rgb(col)
    dia = 0.12
    for m in ms:
        col = TEAL if (ann and m["date"] < ann) else GOLD
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(X(m["date"]) - dia / 2), Inches(axis_y - dia / 2),
                                     Inches(dia), Inches(dia))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(col)
        dot.line.color.rgb = rgb("FFFFFF"); dot.line.width = Pt(0.75)
        dot.shadow.inherit = False
    for d, ha in [(lo, PP_ALIGN.LEFT), (hi, PP_ALIGN.RIGHT)]:
        tb = slide.shapes.add_textbox(Inches(X(d) - (1.4 if ha == PP_ALIGN.RIGHT else 0)),
                                      Inches(axis_y + 0.16), Inches(1.4), Inches(0.22))
        p = tb.text_frame.paragraphs[0]; p.alignment = ha
        r = p.add_run(); r.text = d.strftime("%b %Y"); r.font.size = Pt(8); r.font.color.rgb = rgb(MUT)


def box(slide, l, t, w, h, fill=None, line=None, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = rgb(line); shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        for (s, size, color, bold) in para:
            r = p.add_run(); r.text = s
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = rgb(color); r.font.name = "Calibri"
    return tb


def pct(x):
    x = fnum(x)
    return "—" if x is None else f"{x*100:.0f}%"


def slide_title(s, title, sub):
    box(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=False)
    txt(s, 0.6, 0.45, 12.1, 0.7, [[(title, 30, NAVY, True)]])
    txt(s, 0.62, 1.05, 12.1, 0.4, [[(sub, 14, BLUE, False)]])


def build_deck(rep, deal, series, path):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    name = re.sub(r",? (Inc|Corp|Ltd|Co|Holdings)\.?.*$", "", rep["title"]).strip()
    acq = deal.get("acquirer") or "—"
    dealp = fnum(deal.get("per_share_upfront"))

    # ---------- Slide 1: deal snapshot (dark) ----------
    s = prs.slides.add_slide(blank)
    box(s, 0, 0, 13.333, 7.5, fill=NAVY, radius=False)
    txt(s, 0.7, 0.7, 11.9, 0.5, [[("BUYOUT CASE STUDY", 13, "8FB4E0", True)]])
    txt(s, 0.7, 1.15, 11.9, 1.5,
        [[(name, 40, LIGHT, True)],
         [(f"{rep.get('ticker','')}  ·  acquired by {acq}", 18, "CADCFC", False)]], space=1.05)
    stats = [
        ("Deal price/sh", f"${dealp:g}" if dealp else "—", GOLD),
        ("Premium (unaffected)", pct(deal.get("premium_to_unaffected")), TEAL),
        ("4-wk pre-ann. run-up", pct(deal.get("runup_4wk_pct")), HOT),
        ("% premium pre-captured", pct(deal.get("pct_premium_pre_captured")), HOT),
    ]
    cw, gap, x0 = 2.85, 0.25, 0.7
    for i, (lab, val, col) in enumerate(stats):
        x = x0 + i * (cw + gap)
        box(s, x, 3.1, cw, 1.7, fill=PANEL)
        txt(s, x + 0.2, 3.25, cw - 0.4, 1.0, [[(val, 30, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + 0.2, 4.25, cw - 0.4, 0.5, [[(lab, 12, "9DB2C9", False)]])
    annd = rep["ann"].strftime("%b %d, %Y") if rep["ann"] else "—"
    closed = rep["close"].strftime("%b %d, %Y") if rep["close"] else "n/a"
    vol = deal.get("vol_ratio_4wk")
    foot = (f"Announced {annd}   ·   Closed {closed}   ·   "
            f"Status: {rep.get('status','')}   ·   "
            f"Pre-ann. volume {vol}× normal" if vol else
            f"Announced {annd}   ·   Closed {closed}   ·   Status: {rep.get('status','')}")
    txt(s, 0.7, 5.3, 11.9, 0.5, [[(foot, 12, "8FB4E0", False)]])
    txt(s, 0.7, 6.85, 11.9, 0.4,
        [[("Sources: SEC SC 14D-9 filing + weekly market data. Rumor = price move with no identified filing/news.", 9, "5E7790", False)]])

    # ---------- Slide 2: stock price + rumors + events ----------
    s = prs.slides.add_slide(blank)
    slide_title(s, "Stock price, events & rumors", name)
    add_price_chart(s, rep, deal, series)
    txt(s, 0.6, 6.55, 12.1, 0.7,
        [[("red = unexplained move (rumor/leak)", 11, HOT, True),
          (" · ", 11, MUT, False),
          ("teal = news-driven (event)", 11, TEAL, True),
          (" · ", 11, MUT, False),
          ("gold = deal announcement", 11, GOLD, True)],
         [("Rumor/event markers are placed one week before the reported move week "
           "(weekly data; a reported move can lag the true move by up to a week). "
           "The deal marker is on the actual announcement week.", 9.5, MUT, False)]])

    # ---------- Slide 3+: interaction timeline (diagram + COMPLETE dated list) ----------
    ms = rep["timeline"]
    cap0, capN = 24, 36
    pages, i, first = [], 0, True
    while True:
        cap = cap0 if first else capN
        pages.append(ms[i:i + cap]); i += cap; first = False
        if i >= len(ms):
            break
    for pi, page in enumerate(pages):
        s = prs.slides.add_slide(blank)
        slide_title(s, "Interaction timeline" + ("" if pi == 0 else " (cont.)"), name)
        top = 1.35
        if pi == 0 and ms:
            add_timeline_diagram(s, rep, left=0.7, top=1.4, width=11.9)
            top = 2.85
        rpc = (len(page) + 1) // 2
        xcol = [0.6, 6.85]
        for j, m in enumerate(page):
            col = 0 if j < rpc else 1
            row = j if col == 0 else j - rpc
            yy = top + row * 0.30
            txt(s, xcol[col], yy, 1.45, 0.28, [[(fmt_when(m["date"], m["prec"]), 9, BLUE, True)]])
            txt(s, xcol[col] + 1.5, yy, 4.45, 0.28, [[(clip(m["text"], 60), 9, INK, False)]])
        if pi == 0:
            note = (f"{len(ms)} disclosed milestones, first contact → close · "
                    "dates shown at the precision the filing states (month/year where imprecise).")
            txt(s, 0.6, 7.05, 12.1, 0.4, [[(note, 9, MUT, False)]])

    # ---------- Slide 4: the rumors ----------
    s = prs.slides.add_slide(blank)
    slide_title(s, "The rumors — unexplained pre-announcement moves", name)
    def is_deal_week(w):
        return bool(rep["ann"] and (rep["ann"] - w).days <= 6)
    all_unexpl = [i for i in rep["inflections"]
                  if "not identified" in (i["trigger"] or "").lower() and i["chg"] is not None]
    rumors = [i for i in all_unexpl if not is_deal_week(i["week"])]
    txt(s, 0.6, 1.5, 12.1, 0.5,
        [[(f"{len(rumors)} unexplained pre-announcement move weeks "
           f"(of {len(all_unexpl)} total unexplained inflections)", 15, BLUE, True)]])
    top = sorted([i for i in rumors if i["chg"] is not None],
                 key=lambda i: abs(i["chg"]), reverse=True)[:8]
    y = 2.2
    box(s, 0.6, y, 12.1, 0.5, fill="EAF1F9")
    for cx, w, h in [(0.8, 2.4, "Week ended"), (3.2, 2.6, "Move"), (6.0, 6.5, "Timing vs. announcement")]:
        txt(s, cx, y + 0.06, w, 0.4, [[(h, 12, INK, True)]])
    y += 0.5
    for i in top:
        rel = "—"
        if rep["ann"]:
            dd = (rep["ann"] - i["week"]).days
            rel = (f"{dd} days before announcement" if dd > 0
                   else ("announcement week" if dd > -7 else "after announcement"))
        px = f"  (${i['pre']:g}→${i['post']:g})" if (i["pre"] and i["post"]) else ""
        txt(s, 0.8, y + 0.05, 2.4, 0.4, [[(i["week"].strftime("%Y-%m-%d"), 11, INK, False)]])
        col = HOT if i["chg"] > 0 else MUT
        txt(s, 3.2, y + 0.05, 2.7, 0.4, [[(f"{i['chg']:+.0f}%{px}", 11, col, True)]])
        txt(s, 6.0, y + 0.05, 6.5, 0.4, [[(f"No filing/news identified · {rel}", 11, MUT, False)]])
        y += 0.46
    txt(s, 0.6, 6.85, 12.1, 0.4,
        [[("These are weeks the stock moved materially with no disclosed catalyst — the leak/rumor footprint.",
           10, MUT, False)]])

    # ---------- Slide 5: the events ----------
    s = prs.slides.add_slide(blank)
    slide_title(s, "The events — disclosed catalysts", name)
    cats = [(nd, h, tag_event(h)) for nd, h in rep["news"] if CATALYST_RE.search(h)]
    cats.sort(key=lambda c: (c[0]))
    show = cats[-9:] if len(cats) > 9 else cats
    colmap = {"Deal/M&A": GOLD, "Regulatory": HOT, "Clinical": TEAL, "Partnership": BLUE,
              "Financing": MUT, "Other": MUT}
    y = 1.5
    for nd, h, tg in show:
        box(s, 0.7, y + 0.04, 1.55, 0.34, fill=colmap[tg])
        txt(s, 0.7, y + 0.06, 1.55, 0.3, [[(tg, 10, LIGHT, True)]], align=PP_ALIGN.CENTER)
        txt(s, 2.45, y, 1.5, 0.4, [[(nd.strftime("%Y-%m-%d"), 11, MUT, False)]])
        txt(s, 3.9, y, 8.8, 0.5, [[(clip(h, 108), 11.5, INK, False)]])
        y += 0.56
    if not show:
        txt(s, 0.7, 2.0, 12, 0.5, [[("No catalyst-tagged news in the report.", 13, MUT, False)]])
    txt(s, 0.6, 6.95, 12.1, 0.4,
        [[("Catalysts tagged by keyword from the report's Company News table.", 10, MUT, False)]])

    prs.save(path)


def make_one(arg, deals, prices):
    fp = resolve_file(arg)
    if not fp:
        print(f"  ! not found: {arg}"); return
    rep = parse_report(fp)
    deal = deals.get(fp.name, {})
    if deal.get("data_quality") == "suspect_feed":
        print(f"  ! {fp.name}: suspect price feed — skipping (charts unreliable)"); return
    series = prices.get(fp.name, [])
    if not series:
        print(f"  ! {fp.name}: no price series"); return
    rep["ticker"] = rep.get("ticker") or fp.stem.split("_")[-1].upper()
    out = DECKS / f"{rep['ticker']}.pptx"
    try:
        build_deck(rep, deal, series, out)
    except Exception as e:  # noqa: BLE001
        import traceback
        print(f"  ✗ {fp.name}: {type(e).__name__}: {e}")
        traceback.print_exc()
        return
    print(f"  ✓ {out.name} ({len(rep['timeline'])} milestones)")


def main():
    DECKS.mkdir(parents=True, exist_ok=True)
    args = sys.argv[1:]
    deals, prices = load_dataset()
    if args and args[0] == "--all":
        for fp in sorted(REPORTS.glob("*.md")):
            if fp.name != "README.md":
                make_one(fp.name, deals, prices)
    elif args:
        for a in args:
            make_one(a, deals, prices)
    else:
        print(__doc__)


if __name__ == "__main__":
    main()
